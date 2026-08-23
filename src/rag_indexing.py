#!/usr/bin/env python3
"""
rag_indexing.py — RAG 색인 모듈 [배포 준비 사본 — 원본(rag_demo.py) 기반]

문서를 텍스트로 추출하고 청킹한 뒤, 사전학습된 bge-m3 임베딩 모델로
벡터화해 검색 가능한 상태로 저장하는 파이프라인입니다. ML 모델을
"학습(training)"하는 게 아니라 이미 학습된 모델로 "색인(indexing)"만
수행합니다 — 파일명이 그 점을 정확히 반영합니다.

구조:
  [색인] docs/ 의 문서 → 텍스트 추출 → 청킹 → llama-cpp-python(인프로세스) 임베딩 → FAISS
  [질의] 질문 임베딩 → top-k 검색 → 근거 조각 첨부 → Claude Haiku 답변

v7 변경 (체크포인트 단위를 파일 그룹(100개) → 파일 1개로 축소):
  [CHECKPOINT] v6까지는 RAG_CHECKPOINT_FILES(기본 100)개 파일을 묶어 청크를
      모은 뒤 한 번에 임베딩하고 그룹 단위로 커밋했다. v7부터는 파일 1개를
      추출→청킹→임베딩한 직후 그 파일의 청크만 INSERT하고 즉시 커밋한다 —
      중단/크래시 시 유실 범위가 그룹(최대 100개 파일)에서 파일 1개로
      줄어든다.

v6 변경 (FAISS 전메모리 색인 제거, SQLite+sqlite-vec 디스크 기반 색인 도입):
  [DISKINDEX] 벡터를 FAISS IndexFlatIP(프로세스 메모리에 전체 상주)가 아니라
      SQLite 파일(rag_index.db) + sqlite-vec 확장으로 저장한다.

v5 변경 (Ollama HTTP 제거):
  [EMBED] llama-cpp-python으로 Ollama가 이미 받아둔 bge-m3 GGUF blob을 이
      프로세스 안에서 직접 로드해 임베딩한다.

v4 변경 (노트북 실사용 규모 대응): [FILTER] .ragignore, [CATEGORY] 카테고리
  태깅, [ROBUST] 파일 단위 예외 처리.

이전 버전 기능 유지: [INC-1] 파일별 지문 추적, [INC-2] 임베딩 병렬화,
  [INC-3] 상대경로 source, [DOC-1] 다형식 로더, [DOC-2] 임베딩 견고화.

설정 (.env 또는 settings.json → 환경변수로 주입):
  RAG_DOCS_DIR=<색인할 문서 루트, 절대경로. 콤마(,)로 여러 개 지정 가능>
  RAG_EMBED_WORKERS=4    (embed() 호출 큐잉 동시성 — CPU 코어 점유와는 무관, 아래 참고)
  RAG_EMBED_THREADS=<llama.cpp 내부 스레드 수. 미설정 시 코어수//2(최소 1).
               실제 CPU 점유를 결정하는 값은 이거다 — RAG_EMBED_WORKERS가
               아님(그건 몇 개의 embed() 호출을 동시에 큐잉하느냐일 뿐이고,
               실제 계산은 _MODEL_LOCK으로 항상 직렬화된다)>
  ANTHROPIC_API_KEY=... (MCP 전용 모드에서는 불필요 — rag_serve.py 참고)
  RAG_EMBED_MODEL=bge-m3
  RAG_EMBED_MODEL_PATH=<bge-m3 GGUF blob 경로>
  RAG_DATA_DIR=<색인 DB(rag_index.db) 저장 위치를 강제로 지정(선택, 하위호환용).
               미설정 시 app_paths.get_app_data_dir()가 자동 결정 —
               exe(frozen)면 %APPDATA%/NotebookRAG, 개발 모드면 이 파일과
               같은 디렉터리(src/)>

.ragignore 예시 (각 RAG_DOCS_DIR 루트 바로 아래에 개별 생성):
  # 주석은 # 으로 시작
  COM_ENG                # 폴더명 통째로 제외
  *임시*                 # 파일/폴더명에 "임시" 포함 시 제외
  *.msi
  *최종*_v[0-9]*         # "최종_v2" 같은 중복본 패턴 (필요시 조정)

사용:
  python3 rag_indexing.py               # 증분 색인 후 대화형 질의
  python3 rag_indexing.py --reindex     # 전체 강제 재색인 (모델 교체 등)
"""

from __future__ import annotations

import fnmatch
import gc
import hashlib
import os
import sqlite3
import sys
import threading
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

import numpy as np
from llama_cpp import Llama
from dotenv import load_dotenv

from app_paths import get_app_data_dir, load_settings_json, load_indexer_config, is_frozen, get_install_dir

load_settings_json()  # load_dotenv()보다 먼저 호출 (우선순위: 환경변수 > .env > settings.json)
load_dotenv()

try:
    import sqlite_vec
except ImportError:
    sys.exit("오류: sqlite-vec 패키지가 없습니다 — 디스크 기반 벡터 색인에 "
             "필요합니다. `pip install sqlite-vec` 실행 후 다시 시도하세요.")

# ── 설정 ─────────────────────────────────────────────────────────────────────

BASE_DIR    = Path(__file__).parent
DATA_DIR    = get_app_data_dir()
CACHE_FILE  = DATA_DIR / "rag_index.npz"
CACHE_DB    = DATA_DIR / "rag_index.db"
EMBED_MODEL = os.getenv("RAG_EMBED_MODEL", "bge-m3")
def _default_embed_model_path() -> Path:
    """[티켓 C-1 선결 과제] Ollama blob 기본값은 "개발자 컴퓨터에 이미 Ollama로
    bge-m3를 받아둔 상태"를 전제한 dev 전용 값 — 일반 사용자 배포판(exe)엔
    Ollama 자체가 없으므로 얼어붙은(frozen) 상태에서는 설치 폴더 안의
    models/bge-m3.gguf를 기본값으로 삼는다(모델 자동 다운로드는 별도
    인스톨러 티켓 — 이번 티켓은 그 파일이 수동으로 그 자리에 있다고 가정)."""
    if is_frozen():
        return get_install_dir() / "models" / "bge-m3.gguf"
    return (Path(os.path.expandvars("%USERPROFILE%")) / ".ollama" / "models" / "blobs" /
            "sha256-daec91ffb5dd0c27411bd71f29932917c49cf529a641d0168496c3a501e3062c")


_DEFAULT_EMBED_MODEL_PATH = _default_embed_model_path()
EMBED_MODEL_PATH = Path(os.getenv("RAG_EMBED_MODEL_PATH", str(_DEFAULT_EMBED_MODEL_PATH)))
EMBED_WORKERS = int(os.getenv("RAG_EMBED_WORKERS", "4"))
# [성능 튜닝] RAG_EMBED_WORKERS와는 별개다 — WORKERS는 "몇 개의 embed() 호출을
# 동시에 큐잉하느냐"만 조절하고(실제 계산은 _MODEL_LOCK으로 어차피 직렬화됨),
# CPU 코어 점유는 이 값(llama.cpp 내부 스레드 수)이 결정한다. 예전엔
# os.cpu_count()로 고정돼 있어서 WORKERS를 1로 낮춰도 임베딩 한 번마다 모든
# 코어를 썼다 — NotebookRAG는 상시 백그라운드 서비스라 절반(//2)에서
# 한 번 더 낮춰 코어의 1/4(최소 1)을 기본값으로 삼는다(사용자 실측: //2로도
# 작업관리자 CPU%가 50%대 — 상시 백그라운드 서비스는 그보다 더 존재감이
# 옅어야 한다는 피드백 반영, 목표 25%대).
_DEFAULT_EMBED_THREADS = max(1, (os.cpu_count() or 4) // 4)
EMBED_THREADS = int(os.getenv("RAG_EMBED_THREADS", str(_DEFAULT_EMBED_THREADS)))
if "nomic" in EMBED_MODEL:
    PREFIX_DOC, PREFIX_QUERY = "search_document: ", "search_query: "
else:
    PREFIX_DOC = PREFIX_QUERY = ""
CHUNK_SIZE, CHUNK_OVERLAP = 500, 100
DEGENERATE_COMPRESSION_RATIO = 0.05
DEGENERATE_SAMPLE_CHARS      = 200_000
ABSOLUTE_MAX_CHARS = 3_000_000
EMBED_MAX_RETRY   = 3
TOP_K       = 4
SUPPORTED   = (".md", ".txt", ".pdf", ".docx", ".pptx", ".hwpx", ".hwp")
# [HWP지원] .txt는 extract_text()가 이미 처리하고 있었는데(".txt", ".md" 같은
# 분기) SUPPORTED에는 빠져 있어서 scan_files()의 확장자 필터를 통과하지 못해
# 조용히 무시되고 있었다 — 이번 티켓의 "기존 파일이 조용히 무시되고 있었는지"
# 확인 과정에서 발견(실제 코퍼스에 .txt 2,233개, HWP와 무관한 별개 버그).

# ── [MULTI] 문서 루트 다중 지정 ──────────────────────────────────────────────

def _resolve_docs_dirs(explicit: list[Path] | None = None) -> list[Path]:
    """[티켓 B 선결 과제] scan_files()가 이 함수를 매 회차 다시 호출하므로
    트레이 앱이 프로세스 재시작 없이 추가/삭제한 폴더가 다음 색인 루프
    회차부터 반영된다. 우선순위: 명시적 인자(indexer_serve.py의 색인 루프가
    indexer_config.json 값을 직접 읽어 넘김) > RAG_DOCS_DIR 환경변수(dev/CLI
    하위호환) > indexer_config.json(둘 다 없을 때만) > 기본값(이 파일 옆 docs/).

    [티켓 E] 온보딩 폴더(설치 폴더/onboarding/)는 위 우선순위와 무관하게
    항상 목록 맨 앞에 붙인다 — indexer_config.json(사용자가 /folders로 편집
    가능한 목록)에는 절대 안 들어가므로, 트레이 앱에서 사용자가 자기 폴더를
    전부 지워도 "NotebookRAG 안녕" 온보딩 문서는 계속 검색 대상에 남는다
    (오늘 확정한 "계속 남기자" 결정 — 사용설명서/FAQ 역할까지 겸함).

    [실사용 발견] 원래는 목록 맨 뒤에 붙였는데, 설치 안내 자체가 "모델
    다운로드 → 폴더 추가"를 바로 이어서 하도록 유도하다 보니, 사용자가
    설치 직후 빠르게 큰 폴더를 추가하면 그 폴더가 온보딩과 같은 회차에
    스캔되면서 파일 하나뿐인 온보딩이 대용량 폴더 뒤로 밀려 한참 안 끝나는
    문제가 있었다(실제로 이 개발 PC에서도 대량 재색인 중 재현됨). 맨 앞에
    붙이면 매 회차 온보딩(파일 1개)을 먼저 처리하므로 비용은 거의 안 들면서
    사용자가 폴더를 얼마나 빨리/많이 추가하든 첫 검색 경험이 보장된다."""
    if explicit is not None:
        dirs = explicit
    else:
        raw_env = os.getenv("RAG_DOCS_DIR")
        if raw_env:
            dirs = [Path(p.strip()) for p in raw_env.split(",") if p.strip()]
        else:
            configured = load_indexer_config().get("docs_dirs", [])
            dirs = [Path(p) for p in configured] if configured else [BASE_DIR / "docs"]

    onboarding = get_install_dir() / "onboarding"
    if onboarding.is_dir() and onboarding not in dirs:
        dirs = [onboarding] + dirs
    return dirs

def _make_dir_labels(dirs: list[Path]) -> dict[str, Path]:
    labels: dict[str, Path] = {}
    seen: dict[str, int] = {}
    for d in dirs:
        name = d.name or str(d)
        if name in seen:
            seen[name] += 1
            label = f"{name}({seen[name]})"
        else:
            seen[name] = 1
            label = name
        labels[label] = d
    return labels

DOCS_DIRS  = _resolve_docs_dirs()
DIR_LABELS = _make_dir_labels(DOCS_DIRS)

def resolve_path(rel: str) -> Path:
    label, _, sub = rel.partition("/")
    return DIR_LABELS[label] / sub

# ── [FILTER] .ragignore ──────────────────────────────────────────────────────

def load_ignore_patterns(root: Path) -> list[str]:
    f = root / ".ragignore"
    if not f.exists():
        return []
    return [ln.strip() for ln in f.read_text(encoding="utf-8").splitlines()
            if ln.strip() and not ln.strip().startswith("#")]

def is_ignored(rel: str, patterns: list[str]) -> bool:
    rel_posix = rel.replace("\\", "/")
    parts = rel_posix.split("/")
    for pat in patterns:
        if fnmatch.fnmatch(rel_posix, pat):
            return True
        if any(fnmatch.fnmatch(part, pat) for part in parts):
            return True
    return False

# ── [DOC-1] 형식별 텍스트 추출 ───────────────────────────────────────────────

def extract_text(f: Path) -> str | list[str]:
    suf = f.suffix.lower()
    if suf in (".txt", ".md"):
        return f.read_text(encoding="utf-8", errors="ignore")
    if suf == ".pdf":
        from pypdf import PdfReader
        try:
            return "\n".join(p.extract_text() or "" for p in PdfReader(f).pages)
        except Exception as exc:
            print(f"    경고: {f.name} PDF 추출 실패 — 건너뜀 ({exc})")
            return ""
    if suf == ".docx":
        from docx import Document
        doc = Document(f)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text for c in row.cells))
        return "\n".join(parts)
    if suf == ".pptx":
        from pptx import Presentation
        slides = []
        for n, slide in enumerate(Presentation(f).slides, 1):
            texts = [sh.text for sh in slide.shapes
                     if sh.has_text_frame and sh.text.strip()]
            if texts:
                slides.append(f"[슬라이드 {n}]\n" + "\n".join(texts))
        return slides
    if suf == ".hwpx":
        # [HWP지원] hwpx는 zip+XML(OWPML) 구조라 표준 라이브러리만으로 충분하다
        # — 별도 의존성 없음. Contents/section*.xml 각각의 모든 텍스트 노드를
        # (태그/네임스페이스 무관하게) itertext()로 순서대로 그러모은다.
        import zipfile
        import xml.etree.ElementTree as ET
        try:
            parts = []
            with zipfile.ZipFile(f) as z:
                section_names = sorted(
                    n for n in z.namelist()
                    if n.startswith("Contents/section") and n.endswith(".xml")
                )
                for name in section_names:
                    root = ET.fromstring(z.read(name))
                    text = "".join(root.itertext())
                    if text.strip():
                        parts.append(text)
            return "\n".join(parts)
        except Exception as exc:
            print(f"    경고: {f.name} HWPX 추출 실패 — 건너뜀 ({exc})")
            return ""
    if suf == ".hwp":
        # [HWP지원] 구버전 HWP(복합파일+레코드 구조)는 오픈소스 pyhwp(hwp5)로
        # 충분히 커버됨을 실제 문서 다수로 확인해서, 사용자의 MFC 참조 코드
        # 포팅은 필요 없었다(Go/No-Go 조사 결과: Go, 오픈소스로 충분).
        # hwp5txt CLI와 동일한 변환 파이프라인을 인프로세스로 재사용한다
        # (서브프로세스 없이 — PyInstaller 번들 안에서 그대로 동작).
        import io
        from contextlib import closing
        from hwp5.xmlmodel import Hwp5File
        from hwp5.hwp5txt import TextTransform
        try:
            transform = TextTransform().transform_hwp5_to_text
            buf = io.BytesIO()
            with closing(Hwp5File(str(f))) as hwp5file:
                transform(hwp5file, buf)
            return buf.getvalue().decode("utf-8", errors="ignore")
        except Exception as exc:
            print(f"    경고: {f.name} HWP 추출 실패 — 건너뜀 ({exc})")
            return ""
    return ""

# ── 청킹 ─────────────────────────────────────────────────────────────────────

def chunk_text(text: str, source: str) -> list[dict]:
    chunks = []
    step = CHUNK_SIZE - CHUNK_OVERLAP
    for start in range(0, max(len(text) - CHUNK_OVERLAP, 1), step):
        piece = text[start:start + CHUNK_SIZE].strip()
        if len(piece) > 50:
            chunks.append({"source": source, "text": piece})
    return chunks

def extract_file_chunks(f: Path, rel: str) -> list[dict]:
    rel_posix = rel.replace("\\", "/")
    category = rel_posix.split("/")[0] if "/" in rel_posix else "_root"
    try:
        extracted = extract_text(f)
    except Exception as exc:
        print(f"    경고: {rel} 추출 중 예외 발생 — 건너뜀 ({exc})")
        return []

    if isinstance(extracted, str) and len(extracted) > DEGENERATE_SAMPLE_CHARS:
        import zlib
        sample = extracted[:DEGENERATE_SAMPLE_CHARS].encode("utf-8", errors="ignore")
        ratio = len(zlib.compress(sample, 6)) / max(len(sample), 1)
        if ratio < DEGENERATE_COMPRESSION_RATIO:
            print(f"    ⚠️ 경고: {rel} 반복성 텍스트 의심 (압축률 {ratio:.1%}, 원본 "
                  f"{len(extracted):,}자) — 깨진 추출로 판단, 5,000자만 사용")
            extracted = extracted[:5_000]
        elif len(extracted) > ABSOLUTE_MAX_CHARS:
            print(f"    ⚠️ 경고: {rel} 정상 텍스트이나 매우 큼 (압축률 {ratio:.1%}, "
                  f"{len(extracted):,}자) — 외부 상한 {ABSOLUTE_MAX_CHARS:,}자로 제한")
            extracted = extracted[:ABSOLUTE_MAX_CHARS]

    if isinstance(extracted, list):
        return [{"source": rel, "category": category, "text": s.strip()}
                for s in extracted if len(s.strip()) > 30]
    if extracted.strip():
        chunks = chunk_text(extracted, rel)
        for c in chunks:
            c["category"] = category
        return chunks
    print(f"    경고: {rel} 에서 텍스트를 추출하지 못함 (스캔본 PDF일 수 있음)")
    return []

# ── [INC-1]+[FILTER] 파일별 지문 스캔 ────────────────────────────────────────

def scan_files(docs_dirs: list[Path] | None = None) -> dict[str, tuple[int, int]]:
    """[티켓 B] 매 호출마다 DOCS_DIRS/DIR_LABELS를 다시 계산한다(모듈 임포트
    시점의 상수가 아니라) — 상시 색인 루프에서 폴더 추가/삭제가 재시작 없이
    반영되게 하는 핵심 지점. resolve_path()는 이 호출 직후에만 유효하다는
    기존 호출 순서(scan_files() → 같은 회차 안에서 resolve_path() 사용)를
    그대로 전제한다."""
    global DOCS_DIRS, DIR_LABELS
    DOCS_DIRS = _resolve_docs_dirs(docs_dirs)
    DIR_LABELS = _make_dir_labels(DOCS_DIRS)

    out = {}
    skipped = 0
    any_dir_found = False
    for label, root in DIR_LABELS.items():
        if not root.is_dir():
            print(f"경고: {root}/ 디렉터리가 없습니다 — 건너뜀")
            continue
        any_dir_found = True
        patterns = load_ignore_patterns(root)
        for f in sorted(root.rglob("*")):
            if f.is_file() and f.suffix.lower() in SUPPORTED:
                sub = str(f.relative_to(root))
                if patterns and is_ignored(sub, patterns):
                    skipped += 1
                    continue
                st = f.stat()
                rel = f"{label}/{sub}"
                out[rel] = (st.st_size, st.st_mtime_ns)
    if not any_dir_found:
        sys.exit(f"오류: RAG_DOCS_DIR로 지정된 디렉터리가 하나도 없습니다 "
                 f"({', '.join(str(d) for d in DOCS_DIRS)})")
    if skipped:
        print(f"[색인] .ragignore로 제외된 파일: {skipped}개")
    if not out:
        sys.exit(f"오류: 색인할 문서가 없습니다 (지원 형식: {', '.join(SUPPORTED)})")
    return out

# ── [DEDUP] 콘텐츠 해시 기반 중복 파일 탐지 ──────────────────────────────────

_HASH_CHUNK_BYTES = 1 << 20

def _file_content_hash(f: Path) -> str:
    h = hashlib.sha1()
    with open(f, "rb") as fp:
        for block in iter(lambda: fp.read(_HASH_CHUNK_BYTES), b""):
            h.update(block)
    return h.hexdigest()

def find_content_duplicates(rels: list[str]) -> dict[str, str]:
    dup_of: dict[str, str] = {}
    seen_hashes: dict[str, str] = {}
    for rel in rels:
        try:
            h = _file_content_hash(resolve_path(rel))
        except OSError as exc:
            print(f"    경고: {rel} 해시 계산 실패 — 중복 검사 건너뜀 ({exc})")
            continue
        if h in seen_hashes:
            dup_of[rel] = seen_hashes[h]
        else:
            seen_hashes[h] = rel
    if dup_of:
        print(f"[색인] 콘텐츠 중복 파일 {len(dup_of)}개 발견 (신규/변경 파일 "
              f"{len(rels)}개 중) — 청킹/임베딩 생략:")
        for dup, orig in sorted(dup_of.items()):
            print(f"    중복: {dup}  ==  원본: {orig}")
    return dup_of

# ── [DOC-2]+[INC-2] 임베딩 (견고화 + 병렬, llama-cpp-python 인프로세스) ──────

_MODEL_LOCK = threading.Lock()
_MODEL_LOAD_LOCK = threading.Lock()
_MODEL: Llama | None = None


class ModelNotReadyError(RuntimeError):
    """[티켓 D 선결 과제] 모델 파일이 아직 없음(최초 실행 시 자동 다운로드
    진행 중 등) — 호출부가 이 예외를 잡아 "모델 준비 중"류의 정상적인
    응답(4xx/5xx 중 하나, 크래시 아님)으로 바꿔야 한다. 예전에는 여기서
    sys.exit()을 불렀는데, 상시 서버(notebookrag_main.py) 안에서는 그게
    이 요청 하나만이 아니라 프로세스 전체(모든 API, /health 포함)를
    죽이는 치명적인 문제였다."""


def _get_model() -> Llama:
    global _MODEL
    if _MODEL is not None:
        return _MODEL
    with _MODEL_LOAD_LOCK:
        if _MODEL is None:
            if not EMBED_MODEL_PATH.exists():
                raise ModelNotReadyError(str(EMBED_MODEL_PATH))
            print(f"[임베딩] llama-cpp-python 모델 로드 중: {EMBED_MODEL_PATH}")
            t0 = time.time()
            # [버그 수정 — 사용자 실측(95% CPU, 스레드 99개)으로 발견] n_threads만
            # 지정하고 n_threads_batch를 안 주면 llama-cpp-python이
            # n_threads_batch를 자기 내부 기본값(multiprocessing.cpu_count() —
            # 전체 코어)으로 따로 정한다(llama_cpp/llama.py:
            # "self.n_threads_batch = n_threads_batch or multiprocessing.cpu_count()").
            # 임베딩(텍스트 조각 하나를 벡터화)은 llama.cpp 내부적으로 배치
            # 연산(batch/prompt eval)으로 처리되므로, 실제 CPU 점유를 결정하는
            # 건 n_threads(생성용)가 아니라 n_threads_batch다 — EMBED_THREADS를
            # 코어의 절반으로 낮췄는데도 실제로는 전체 코어를 계속 쓰고 있던
            # 원인이 이거였다. 둘 다 EMBED_THREADS로 맞춘다.
            _MODEL = Llama(
                model_path=str(EMBED_MODEL_PATH),
                embedding=True,
                n_threads=EMBED_THREADS,
                n_threads_batch=EMBED_THREADS,
                verbose=False,
            )
            print(f"[임베딩] 모델 로드 완료 ({time.time() - t0:.1f}초)")
    return _MODEL


def _embed_one(text: str) -> list[float] | None:
    text = "".join(ch for ch in text if ch in "\n\t" or ord(ch) >= 32)
    if not text.strip():
        text = " "
    model = _get_model()
    for attempt in range(1, EMBED_MAX_RETRY + 1):
        try:
            with _MODEL_LOCK:
                vec = model.embed(text)
            if isinstance(vec[0], list):
                vec = np.mean(np.array(vec, dtype=np.float32), axis=0).tolist()
            return vec
        except Exception as exc:
            if attempt == EMBED_MAX_RETRY:
                print(f"    ⚠️ 임베딩 실패(재시도 {EMBED_MAX_RETRY}회 소진): {exc}")
                return None
            time.sleep(0.5)
    return None


def _eta_seconds(done: int, total: int, elapsed: float) -> float:
    """indexer_state.IndexerState.to_dict()도 동일한 공식을 재사용한다 —
    ETA 계산 로직을 두 곳에 따로 두지 않기 위함."""
    rate = done / elapsed if elapsed > 0 else 0
    return (total - done) / rate if rate > 0 else 0


def _print_progress(label: str, done: int, total: int, t0: float) -> None:
    elapsed = time.time() - t0
    eta_sec = _eta_seconds(done, total, elapsed)
    print(f"    {label}: {done:,}/{total:,} ({done/total:.1%}) — "
          f"{elapsed/60:.1f}분 경과, 예상 잔여 약 {eta_sec/60:.1f}분")


def _should_report(last_print: list, done: int, total: int, min_interval: float = 3.0) -> bool:
    now = time.time()
    if done == total or now - last_print[0] >= min_interval:
        last_print[0] = now
        return True
    return False


def embed(texts: list[str], label: str = "", workers: int = 1) -> np.ndarray:
    if not texts:
        return np.zeros((0, 0), dtype=np.float32)
    t0 = time.time()
    n = len(texts)
    vecs: list = [None] * n

    if workers <= 1 or n <= 1:
        last_print = [t0]
        for i, t in enumerate(texts):
            vecs[i] = _embed_one(t)
            if label and _should_report(last_print, i + 1, n):
                _print_progress(label, i + 1, n, t0)
    else:
        last_print = [t0]
        done = 0
        with ThreadPoolExecutor(max_workers=workers) as ex:
            futures = {ex.submit(_embed_one, t): i for i, t in enumerate(texts)}
            for fut in as_completed(futures):
                i = futures[fut]
                vecs[i] = fut.result()
                done += 1
                if label and _should_report(last_print, done, n):
                    _print_progress(label, done, n, t0)

    if any(v is None for v in vecs):
        raise RuntimeError("embed(): 일부 텍스트 임베딩 실패(타임아웃) — "
                           "대량 문서처리에는 embed_chunks() 사용")
    m = np.array(vecs, dtype=np.float32)
    return m / (np.linalg.norm(m, axis=1, keepdims=True) + 1e-12)


def embed_chunks(chunks: list[dict], workers: int = 1,
                 label: str = "", on_progress=None) -> tuple[list[dict], np.ndarray]:
    """on_progress(done, total)는 label과 무관하게(레이블이 없어도) 호출된다 —
    [버그 수정] 예전엔 상시 색인 루프가 파일 하나를 embed_chunks()에 넘긴 뒤
    끝날 때까지 진행률(state.set_progress)이 (0, N)에 고정돼서, 대용량
    파일(청크 수백~수천 개) 처리 중엔 /indexer/status가 몇 분씩 멈춰있는
    것처럼 보였다(사용자 보고: "청킹 작업이 멈춰 있음"). _should_report와
    같은 3초 간격 게이트를 공유해 락 경합 없이 주기적으로 갱신한다."""
    if not chunks:
        return [], np.zeros((0, 0), dtype=np.float32)
    texts = [PREFIX_DOC + c["text"] for c in chunks]
    t0 = time.time()
    n = len(texts)
    results: list = [None] * n

    if workers <= 1 or n <= 1:
        last_print = [t0]
        for i, t in enumerate(texts):
            results[i] = _embed_one(t)
            if _should_report(last_print, i + 1, n):
                if label:
                    _print_progress(label, i + 1, n, t0)
                if on_progress:
                    on_progress(i + 1, n)
    else:
        last_print = [t0]
        done = 0
        with ThreadPoolExecutor(max_workers=workers) as ex:
            futures = {ex.submit(_embed_one, t): i for i, t in enumerate(texts)}
            for fut in as_completed(futures):
                i = futures[fut]
                results[i] = fut.result()
                done += 1
                if _should_report(last_print, done, n):
                    if label:
                        _print_progress(label, done, n, t0)
                    if on_progress:
                        on_progress(done, n)

    ok_chunks, ok_vecs = [], []
    for c, v in zip(chunks, results):
        if v is None:
            preview = c["text"][:40].replace("\n", " ")
            print(f"    ⚠️ 경고: {c['source']} 조각 임베딩 실패(타임아웃, 깨진 "
                  f"디코딩 의심) — 건너뜀: {preview!r}...")
        else:
            ok_chunks.append(c)
            ok_vecs.append(v)

    if not ok_vecs:
        return [], np.zeros((0, 0), dtype=np.float32)
    m = np.array(ok_vecs, dtype=np.float32)
    m = m / (np.linalg.norm(m, axis=1, keepdims=True) + 1e-12)
    return ok_chunks, m

# ── [DISKINDEX] SQLite+sqlite-vec 저장소 ─────────────────────────────────────

_DB_LOCK = threading.Lock()
_SQL_BATCH = 500

def _gen_b_path(primary: Path) -> Path:
    return primary.parent / f"{primary.stem}_b{primary.suffix}"

def _active_pointer_path(primary: Path) -> Path:
    return primary.parent / f"{primary.stem}.active"

def _read_active_gen(primary: Path) -> str:
    p = _active_pointer_path(primary)
    if p.exists():
        v = p.read_text(encoding="utf-8").strip()
        if v in ("a", "b"):
            return v
    return "a"

def _write_active_gen(primary: Path, gen: str) -> None:
    p = _active_pointer_path(primary)
    tmp = p.with_suffix(".tmp")
    tmp.write_text(gen, encoding="utf-8")
    tmp.replace(p)

def _db_path_for_gen(primary: Path, gen: str) -> Path:
    return primary if gen == "a" else _gen_b_path(primary)

def _open_db(path: Path) -> sqlite3.Connection:
    conn = sqlite3.connect(str(path), check_same_thread=False)
    conn.enable_load_extension(True)
    sqlite_vec.load(conn)
    conn.enable_load_extension(False)
    conn.execute("PRAGMA journal_mode=WAL")
    return conn

def _table_exists(conn: sqlite3.Connection, name: str) -> bool:
    row = conn.execute(
        "SELECT name FROM sqlite_master WHERE type IN ('table','view') AND name=?",
        (name,)).fetchone()
    return row is not None

def _ensure_base_schema(conn: sqlite3.Connection) -> None:
    conn.execute("CREATE TABLE IF NOT EXISTS meta (key TEXT PRIMARY KEY, value TEXT)")
    conn.execute("CREATE TABLE IF NOT EXISTS file_meta "
                 "(rel TEXT PRIMARY KEY, size INTEGER, mtime_ns INTEGER)")
    conn.execute("CREATE TABLE IF NOT EXISTS chunks "
                 "(rowid INTEGER PRIMARY KEY, source TEXT, category TEXT, text TEXT)")
    conn.commit()

def _ensure_vec_table(conn: sqlite3.Connection, dim: int) -> None:
    conn.execute(f"CREATE VIRTUAL TABLE IF NOT EXISTS vec_items USING vec0(embedding float[{dim}])")
    conn.commit()

def _drop_all_tables(conn: sqlite3.Connection) -> None:
    for t in ("vec_items", "chunks", "file_meta", "meta"):
        conn.execute(f"DROP TABLE IF EXISTS {t}")
    conn.commit()

def _read_meta(conn: sqlite3.Connection) -> dict[str, str]:
    return dict(conn.execute("SELECT key, value FROM meta").fetchall())

def _read_file_meta(conn: sqlite3.Connection) -> dict[str, tuple[int, int]]:
    return {rel: (size, mtime_ns) for rel, size, mtime_ns in
            conn.execute("SELECT rel, size, mtime_ns FROM file_meta").fetchall()}

def _next_rowid(conn: sqlite3.Connection) -> int:
    row = conn.execute("SELECT COALESCE(MAX(rowid), -1) FROM chunks").fetchone()
    return row[0] + 1

def _delete_sources(conn: sqlite3.Connection, sources: list[str]) -> None:
    if not sources:
        return
    has_vec = _table_exists(conn, "vec_items")
    for i in range(0, len(sources), _SQL_BATCH):
        batch = sources[i:i + _SQL_BATCH]
        placeholders = ",".join("?" * len(batch))
        if has_vec:
            rowids = [r[0] for r in conn.execute(
                f"SELECT rowid FROM chunks WHERE source IN ({placeholders})", batch).fetchall()]
            if rowids:
                conn.executemany("DELETE FROM vec_items WHERE rowid=?", [(r,) for r in rowids])
        conn.execute(f"DELETE FROM chunks WHERE source IN ({placeholders})", batch)
        conn.execute(f"DELETE FROM file_meta WHERE rel IN ({placeholders})", batch)
    conn.commit()

def _migrate_npz_if_needed(conn: sqlite3.Connection, db_path: Path) -> None:
    if not CACHE_FILE.exists():
        return
    if _read_meta(conn):
        return
    z = np.load(CACHE_FILE, allow_pickle=True)
    old_chunks = list(z["chunks"])
    old_matrix = z["matrix"]
    old_files = dict(z["file_meta"].item())
    old_model = str(z.get("embed_model", ""))
    if not old_chunks or old_matrix.size == 0:
        return
    dim = old_matrix.shape[1]
    print(f"[마이그레이션] 기존 {CACHE_FILE.name} 발견 → {db_path.name}"
          f"(SQLite+sqlite-vec)로 1회 이전 중... (조각 {len(old_chunks)}개, 차원 {dim})")
    _ensure_vec_table(conn, dim)
    with conn:
        conn.execute("INSERT OR REPLACE INTO meta VALUES ('embed_model', ?)", (old_model,))
        conn.execute("INSERT OR REPLACE INTO meta VALUES ('dim', ?)", (str(dim),))
        conn.executemany("INSERT OR REPLACE INTO file_meta VALUES (?,?,?)",
                          [(rel, fp[0], fp[1]) for rel, fp in old_files.items()])
        for i, c in enumerate(old_chunks):
            conn.execute("INSERT INTO chunks (rowid, source, category, text) VALUES (?,?,?,?)",
                         (i, c["source"], c.get("category", "_root"), c["text"]))
            conn.execute("INSERT INTO vec_items (rowid, embedding) VALUES (?,?)",
                         (i, np.asarray(old_matrix[i], dtype=np.float32).tobytes()))
    print(f"[마이그레이션] 완료 — {CACHE_FILE.name}은 삭제하지 않고 그대로 "
          f"보존합니다(더 이상 읽지는 않음)")


class ChunkStore:
    def __init__(self, conn: sqlite3.Connection, count: int):
        self._conn = conn
        self._count = count

    def __len__(self) -> int:
        return self._count

    def __getitem__(self, rowid) -> dict:
        with _DB_LOCK:
            row = self._conn.execute(
                "SELECT source, category, text FROM chunks WHERE rowid=?",
                (int(rowid),)).fetchone()
        if row is None:
            raise KeyError(rowid)
        return {"source": row[0], "category": row[1], "text": row[2]}


def _finalize_index(conn: sqlite3.Connection, total_count: int,
                     db_path: Path, gen: str):
    chunks = ChunkStore(conn, total_count)

    def search(qv, k):
        if total_count == 0:
            return np.array([], dtype=np.int64), np.array([], dtype=np.float32)
        kk = min(k, total_count)
        with _DB_LOCK:
            rows = conn.execute(
                "SELECT rowid, distance FROM vec_items WHERE embedding MATCH ? AND k = ? "
                "ORDER BY distance",
                (np.asarray(qv, dtype=np.float32).tobytes(), kk)).fetchall()
        idx = np.array([r[0] for r in rows], dtype=np.int64)
        dist = np.array([r[1] for r in rows], dtype=np.float32)
        scores = 1.0 - (dist * dist) / 2.0
        return idx, scores

    print(f"[색인] sqlite-vec 적재: 벡터 {total_count}개 "
          f"(임베딩: llama-cpp-python 인프로세스, {EMBED_MODEL_PATH.name}) — "
          f"디스크 기반 저장 ({db_path.name}, 세대 {gen})")

    return chunks, search


def get_embed_dim() -> int:
    """[상태정보확장] 활성 색인 DB의 meta 테이블에서 임베딩 차원만 가볍게
    조회한다 — RagRA 인스턴스를 안 거치고 별도 커넥션을 열었다가 바로
    닫으므로 /model/status 같은 자주 폴링되는 엔드포인트에서 불러도 부담이
    적다."""
    active_gen = _read_active_gen(CACHE_DB)
    db_path = _db_path_for_gen(CACHE_DB, active_gen)
    if not db_path.exists():
        return 0
    with _DB_LOCK:
        conn = _open_db(db_path)
        try:
            _ensure_base_schema(conn)
            meta = _read_meta(conn)
        finally:
            conn.close()
    return int(meta["dim"]) if "dim" in meta else 0


def get_file_meta_count() -> int:
    """[DB저장파일수] 활성 색인 DB의 file_meta 테이블 행 수(= DB에 영속
    저장된 고유 파일 개수)를 가볍게 조회한다 — get_embed_dim()과 동일한
    패턴(별도 커넥션 열고 바로 닫음). 감시 폴더의 실제 파일 개수
    (scan_files() 결과, IndexerState.디렉토리총파일수)와 나란히 비교하면
    "색인 안 됐거나 실패한 파일이 있는지"를 한눈에 알 수 있다."""
    active_gen = _read_active_gen(CACHE_DB)
    db_path = _db_path_for_gen(CACHE_DB, active_gen)
    if not db_path.exists():
        return 0
    with _DB_LOCK:
        conn = _open_db(db_path)
        try:
            _ensure_base_schema(conn)
            count = conn.execute("SELECT COUNT(*) FROM file_meta").fetchone()[0]
        finally:
            conn.close()
    return int(count)


def open_existing_index():
    active_gen = _read_active_gen(CACHE_DB)
    db_path = _db_path_for_gen(CACHE_DB, active_gen)

    if not db_path.exists():
        print(f"[색인] 경고: 색인이 없습니다({db_path.name}) — "
              f"rag_indexing.py --reindex를 먼저 실행하거나 RAG.REINDEX 신호를 보내세요")
        return ChunkStore(None, 0), lambda qv, k: (
            np.array([], dtype=np.int64), np.array([], dtype=np.float32))

    with _DB_LOCK:
        conn = _open_db(db_path)
        _ensure_base_schema(conn)
        total_count = conn.execute("SELECT COUNT(*) FROM chunks").fetchone()[0]

    return _finalize_index(conn, total_count, db_path, active_gen)


# ── [INC-1] 증분 색인 ────────────────────────────────────────────────────────

def build_index(force: bool = False, docs_dirs: list[Path] | None = None,
                 state=None, running: threading.Event | None = None):
    """[티켓 B] state(IndexerState)/running(threading.Event)은 상시 색인
    루프(indexer_serve.py)에서만 넘겨준다 — 둘 다 optional이라 CLI 단독
    실행(`python rag_indexing.py --reindex`)은 기존과 동일하게 동작한다."""
    if state is not None:
        state.set_phase("scanning")
    current = scan_files(docs_dirs)
    if state is not None:
        # [DB저장파일수 비교용] phase(idle/processing 등)와 무관하게 항상
        # 최신 값을 유지해야 하므로 start_round()의 _reset_locked()가 안
        # 건드리는 별도 필드에 저장한다(회차 진행 상태와 독립적인 상시 정보).
        state.set_dir_total(len(current))

    with _DB_LOCK:
        active_gen = _read_active_gen(CACHE_DB)
        if force:
            target_gen = "b" if active_gen == "a" else "a"
        else:
            target_gen = active_gen
        db_path = _db_path_for_gen(CACHE_DB, target_gen)

        conn = _open_db(db_path)
        _ensure_base_schema(conn)

        if force:
            _drop_all_tables(conn)
            _ensure_base_schema(conn)
        else:
            _migrate_npz_if_needed(conn, db_path)

        meta = _read_meta(conn)
        old_model = meta.get("embed_model", "")
        dim = int(meta["dim"]) if "dim" in meta else 0
        old_files = _read_file_meta(conn)

        # [긴급 안전장치 — 2026-08-21 실제 데이터 유실 사고 재발 방지] 스캔
        # 결과가 완전히 비어있는데 기존에 색인된 파일이 있었다면, 정상적인
        # "문서가 다 사라짐"이 아니라 감시 폴더 설정을 잘못 읽었거나(경쟁
        # 상태 등) 폴더가 일시적으로 접근 불가능한 상황일 가능성이 훨씬
        # 크다. 이 경우 아래 removed/_delete_sources 로직이 기존 임베딩
        # 전체(수백 MB)를 지워버리는 실제 사고가 있었다 — 통신 구조를
        # HTTP로 바꿔도 원인이 다른 경로(폴더 일시 접근 불가 등)일 수 있어
        # 이중 방어로 계속 유지한다. force=True(명시적 전체 재색인 요청)는
        # 사용자가 의도한 것이므로 이 안전장치를 적용하지 않는다.
        if not force and not current and old_files:
            log_msg = f"스캔 결과가 완전히 비어있음(기존 색인 {len(old_files)}건) — 삭제를 건너뜁니다"
            print(f"[색인] ⚠️ 경고: {log_msg}")
            if state is not None:
                state.add_warning("스캔 결과 비정상 — 삭제 안전장치 발동, 이번 회차 건너뜀")
            total_count = conn.execute("SELECT COUNT(*) FROM chunks").fetchone()[0]
            return _finalize_index(conn, total_count, db_path, target_gen)

        model_changed = old_model != EMBED_MODEL
        if model_changed and old_files:
            print(f"[색인] 임베딩 모델 변경 감지({old_model or '없음'} → "
                  f"{EMBED_MODEL}) — 전체 재색인")
        if model_changed:
            _drop_all_tables(conn)
            _ensure_base_schema(conn)
            old_files = {}
            dim = 0

        reuse_files = {} if (force or model_changed) else {
            rel: fp for rel, fp in current.items()
            if old_files.get(rel) == fp
        }
        changed_files = [rel for rel in current if rel not in reuse_files]
        removed = list(set(old_files) - set(current))

        if removed:
            print(f"[색인] 삭제/제외 감지: {len(removed)}개 파일 — 인덱스에서 제외")
            _delete_sources(conn, removed)

        if changed_files:
            dup_of = find_content_duplicates(changed_files)
            files_to_embed = [rel for rel in changed_files if rel not in dup_of]

            dedup_note = f", 콘텐츠 중복 {len(dup_of)}개 제외" if dup_of else ""
            print(f"[색인] 신규/변경 {len(changed_files)}개{dedup_note} "
                  f"→ 실제 처리 {len(files_to_embed)}개, 재사용 {len(reuse_files)}개 "
                  f"(전체 {len(current)}개 문서)")

            _delete_sources(conn, changed_files)

            warmed_up = False
            total_files = len(files_to_embed)
            total_new_chunks = 0
            cats: dict[str, int] = {}

            if state is not None:
                state.start_round(total_files, len(reuse_files), len(dup_of))

            for i, rel in enumerate(files_to_embed, 1):
                if running is not None and not running.is_set():
                    if state is not None:
                        state.set_phase("paused")
                    running.wait()  # 일시정지는 다음 파일 처리 시작 시점에 반영
                    if state is not None:
                        state.set_phase("processing")

                fname = Path(rel).name
                print(f"[색인] [{i:,}/{total_files:,}] {fname} 처리 중")

                # [메모리 누수 완화 — 2026-08-21 사용자와 확인] python-docx/
                # python-pptx가 내부적으로 쓰는 lxml 트리는 부모-자식 양방향
                # 참조라 순환 참조가 잘 생기는데, 파이썬의 순환 참조 수거기가
                # 자동 임계값 기준으로만 돌다 보니 파일을 빠르게 계속 열고
                # 닫으면 수거 속도가 처리 속도를 못 따라가서 프로세스 메모리가
                # 계속 쌓이는 게 실측됨(파일당 약 7~8MB, 정체 없이 선형 증가).
                # 동작을 바꾸는 게 아니라 정리를 더 자주 강제하는 것뿐이라
                # 안전하다.
                if i % 20 == 0:
                    gc.collect()

                try:
                    if state is not None:
                        state.set_progress(fname, "extract")

                    file_chunks = extract_file_chunks(resolve_path(rel), rel)
                    ok_chunks, file_matrix = [], np.zeros((0, 0), dtype=np.float32)
                    if file_chunks:
                        if not warmed_up:
                            print("[색인] 임베딩 모델 워밍업 시작...")
                            embed(["워밍업"], label="워밍업")
                            print("[색인] 워밍업 완료 — 임베딩 모델 응답 확인됨")
                            warmed_up = True
                        if state is not None:
                            state.set_progress(fname, "embed", (0, len(file_chunks)))
                        on_progress = (
                            (lambda done, total: state.set_progress(fname, "embed", (done, total)))
                            if state is not None else None)
                        ok_chunks, file_matrix = embed_chunks(
                            file_chunks, workers=EMBED_WORKERS, label=f"    {fname}",
                            on_progress=on_progress)

                    with conn:
                        if ok_chunks:
                            if dim == 0:
                                dim = file_matrix.shape[1]
                                _ensure_vec_table(conn, dim)
                                conn.execute("INSERT OR REPLACE INTO meta VALUES ('embed_model', ?)", (EMBED_MODEL,))
                                conn.execute("INSERT OR REPLACE INTO meta VALUES ('dim', ?)", (str(dim),))
                            rid0 = _next_rowid(conn)
                            for j, c in enumerate(ok_chunks):
                                rid = rid0 + j
                                conn.execute(
                                    "INSERT INTO chunks (rowid, source, category, text) VALUES (?,?,?,?)",
                                    (rid, c["source"], c.get("category", "_root"), c["text"]))
                                conn.execute(
                                    "INSERT INTO vec_items (rowid, embedding) VALUES (?,?)",
                                    (rid, np.asarray(file_matrix[j], dtype=np.float32).tobytes()))
                            total_new_chunks += len(ok_chunks)
                            for c in ok_chunks:
                                cat = c.get("category", "_root")
                                cats[cat] = cats.get(cat, 0) + 1
                        fp = current[rel]
                        conn.execute("INSERT OR REPLACE INTO file_meta VALUES (?,?,?)",
                                     (rel, fp[0], fp[1]))

                    if state is not None:
                        if not ok_chunks:
                            kind = "처리실패"
                            state.add_warning(f"{fname}: 처리 실패(추출/임베딩 결과 없음)")
                        else:
                            kind = "변경" if rel in old_files else "신규"
                        state.file_done(len(ok_chunks), kind)
                except Exception as exc:
                    # [버그 수정 — 예외 루틴 누락] 파일 하나가 던진 예외(손상된
                    # 문서, 인코딩 오류 등)가 예전엔 이 for 루프를 그대로 뚫고
                    # build_index() 전체를 죽였다 — 그러면 이미 위에서
                    # _delete_sources()로 지워둔 이번 회차 대상 파일들(이 파일
                    # 포함, 아직 처리 못 한 나머지도)이 전부 미처리 상태로
                    # 남아 다음 회차에 처음부터 다시 시도되고, 같은 파일이
                    # 계속 같은 이유로 죽으면 회차가 영원히 못 끝나는 구조였다.
                    # 이 파일만 실패로 기록하고 다음 파일로 계속 진행한다 —
                    # file_meta도 기록해서(다른 실패 케이스와 동일) 다음 회차에
                    # 무한 재시도하지 않게 한다(내용이 바뀌면 fp가 달라져서
                    # 자연히 다시 시도됨).
                    print(f"    ⚠️ 경고: {fname} 처리 중 예외 발생 — 건너뜀: {exc}")
                    try:
                        with conn:
                            fp = current[rel]
                            conn.execute("INSERT OR REPLACE INTO file_meta VALUES (?,?,?)",
                                         (rel, fp[0], fp[1]))
                    except Exception:
                        pass
                    if state is not None:
                        state.add_warning(f"{fname}: 처리 중 예외 — 건너뜀 ({exc})")
                        state.file_done(0, "처리실패")
                    continue

            if dup_of:
                with conn:
                    for rel in dup_of:
                        fp = current[rel]
                        conn.execute("INSERT OR REPLACE INTO file_meta VALUES (?,?,?)",
                                     (rel, fp[0], fp[1]))

            if total_new_chunks:
                print(f"[색인] 전체 완료 (신규/갱신 조각 {total_new_chunks:,}개, 차원 {dim})")
                print("[색인] 카테고리별 조각 수(이번 처리분): " +
                      ", ".join(f"{k}={v}" for k, v in sorted(cats.items(), key=lambda x: -x[1])[:10]))
        else:
            print(f"[색인] 변경 없음 — 캐시 그대로 사용")
            if state is not None:
                state.start_round(0, len(reuse_files), 0)

        with conn:
            if dim:
                conn.execute("INSERT OR REPLACE INTO meta VALUES ('embed_model', ?)", (EMBED_MODEL,))
                conn.execute("INSERT OR REPLACE INTO meta VALUES ('dim', ?)", (str(dim),))

        total_count = conn.execute("SELECT COUNT(*) FROM chunks").fetchone()[0]

        if force:
            _write_active_gen(CACHE_DB, target_gen)

    return _finalize_index(conn, total_count, db_path, target_gen)

# ── 질의 루프 (CLI 단독 실행용) ───────────────────────────────────────────────

RAG_PROMPT = """당신은 문서 기반 질의응답 도우미입니다.
아래 [근거 문서] 조각들만 사용해 질문에 답하세요.

규칙:
- 근거에 없는 내용은 추측하지 말고 "제공된 문서에서 찾을 수 없습니다"라고 답하세요.
- 답변에 사용한 근거의 번호를 문장 끝에 [1] 형식으로 표기하세요.
- 간결하게 답하세요.

[근거 문서]
{context}

[질문]
{question}"""

def main():
    try:
        chunks, search = build_index(force="--reindex" in sys.argv)
    except ModelNotReadyError as exc:
        sys.exit(f"오류: 임베딩 모델 GGUF 파일을 찾을 수 없습니다: {exc}\n"
                 f"  .env의 RAG_EMBED_MODEL_PATH를 확인하거나, "
                 f"Ollama가 bge-m3를 받아뒀는지(`ollama list`) 확인하세요.")

    from langchain_anthropic import ChatAnthropic
    llm = ChatAnthropic(model="claude-haiku-4-5", temperature=0.0, max_tokens=512)

    print("\nRAG 질의 시작 (종료: 빈 줄 입력)")
    print("─" * 60)
    while True:
        try:
            q = input("\n질문> ").strip()
        except (EOFError, KeyboardInterrupt):
            break
        if not q:
            break

        qv = embed([PREFIX_QUERY + q])[0]
        idx, scores = search(qv, TOP_K)

        if len(idx) == 0:
            print("\n검색 결과가 없습니다 — 색인된 문서가 없습니다.")
            continue

        context = "\n\n".join(
            f"[{n + 1}] (출처: {chunks[int(i)]['source']})\n{chunks[int(i)]['text']}"
            for n, i in enumerate(idx))

        answer = llm.invoke(RAG_PROMPT.format(context=context, question=q))
        print("\n" + answer.content)

        print("\n[근거]")
        for n, (i, sc) in enumerate(zip(idx, scores)):
            c = chunks[int(i)]
            preview = c["text"][:60].replace("\n", " ")
            print(f"  [{n + 1}] [{c.get('category','_root')}] {c['source']} "
                  f"(유사도 {sc:.3f}) {preview}...")


if __name__ == "__main__":
    main()
